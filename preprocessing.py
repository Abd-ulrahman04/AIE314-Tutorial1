import os
import json
import fitz  # PDF
import docx
import pandas as pd
from pptx import Presentation


# ==============================
# Text Cleaning Function
# ==============================
def clean_text(text):
    if not text:
        return ""

    # إزالة المسافات الزائدة
    text = text.replace("\n", " ").replace("\r", " ")
    text = " ".join(text.split())

    return text


# ==============================
# PDF Extraction
# ==============================
def extract_pdf_text(file_path):
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return text


# ==============================
# Word Extraction
# ==============================
def extract_docx_text(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
    return text


# ==============================
# Excel Extraction
# ==============================
def extract_excel_text(file_path):
    text = ""
    try:
        df = pd.read_excel(file_path)
        text = df.to_string()
    except Exception as e:
        print(f"Error reading Excel {file_path}: {e}")
    return text


# ==============================
# PowerPoint Extraction
# ==============================
def extract_ppt_text(file_path):
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPT {file_path}: {e}")
    return text


# ==============================
# Save Normalized JSON
# ==============================
def save_json(title, content, source_file, output_path):
    data = {
        "title": title,
        "content": clean_text(content),
        "metadata": {
            "source_file": source_file,
            "file_type": source_file.split(".")[-1]
        }
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)


# ==============================
# Main Processing Function
# ==============================
def process_files():
    input_folder = "data"
    output_folder = "output"

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for file in os.listdir(input_folder):
        file_path = os.path.join(input_folder, file)

        if not os.path.isfile(file_path):
            continue

        print(f"Processing: {file}")

        title = os.path.splitext(file)[0]
        content = ""

        if file.endswith(".pdf"):
            content = extract_pdf_text(file_path)

        elif file.endswith(".docx"):
            content = extract_docx_text(file_path)

        elif file.endswith(".xlsx"):
            content = extract_excel_text(file_path)

        elif file.endswith(".pptx"):
            content = extract_ppt_text(file_path)

        else:
            print(f"Unsupported file type: {file}")
            continue

        output_path = os.path.join(output_folder, title + ".json")
        save_json(title, content, file, output_path)

    print("Processing completed!")



if __name__ == "__main__":
    process_files()
